"""
RetroTutor Cloud — Retroalimentación automática de evidencias estudiantiles.
Versión desplegable en Streamlit Cloud. Soporta documentos, texto e imágenes.
"""

import base64
import hashlib
import hmac
import io
import os
import tempfile
import zipfile
from datetime import date, datetime, timezone
from pathlib import Path

import pandas as pd
import streamlit as st
from openai import OpenAI

# ── Configuración ─────────────────────────────────────────────────────────────

st.set_page_config(page_title="RetroTutor", page_icon="🎓", layout="wide")

EXT_DOCUMENTO = ["pdf", "docx", "pptx", "txt", "md"]
EXT_IMAGEN = ["png", "jpg", "jpeg", "gif", "webp", "bmp"]
TODAS_EXT = EXT_DOCUMENTO + EXT_IMAGEN

PROMPT = """You are a university tutor. Write SHORT personalized feedback (MAX 150 words, this is critical — never exceed it).

Student: {nombre}
{bloque_guia}
{bloque_rubrica}
VERDICT: {veredicto_texto}
{bloque_contenido}
{bloque_observaciones}

Structure (no headers, no bullets, plain text only):
- Greet by first name
- ONE paragraph: what they did well/poorly, be specific about their content
- TWO concrete suggestions (most impactful)
- Brief encouraging closing

STRICT LIMIT: 150 words maximum. Be concise."""

VEREDICTO_APROBADO = (
    "APPROVED — the student met the requirements. Focus on strengths, "
    "celebrate what they did well, and offer minor suggestions for excellence."
)
VEREDICTO_DEFICIENTE = (
    "NEEDS IMPROVEMENT — the student did not meet the requirements. "
    "Be constructive: clearly state what fell short, give specific actionable "
    "steps to improve, and maintain an encouraging tone. Motivate them to try again."
)

# ── Licencia (compatible con Zajuna Panel) ───────────────────────────────────

ZP_SECRET = "ftM_k2tzaEmrf7hBFE_xTarWJJDfSELX"
ZP_EPOCH = datetime(2026, 1, 1, tzinfo=timezone.utc)


def zp_sign(payload: str) -> str:
    sig = hmac.new(ZP_SECRET.encode(), payload.encode(), hashlib.sha256).hexdigest()
    return sig[:6].upper()


def zp_validate(code: str) -> dict:
    code = (code or "").strip().upper()
    import re
    m = re.match(r"^ZP-([0-9A-Z]{3})-([0-9A-Z]{4})-([0-9A-Z]{6})$", code)
    if not m:
        return {"valid": False, "reason": "Formato inválido. Ejemplo: ZP-XXX-XXXX-XXXXXX"}
    exp, rand, sig = m.group(1), m.group(2), m.group(3)
    expected = zp_sign(f"{exp}-{rand}")
    if sig != expected:
        return {"valid": False, "reason": "Código no válido."}
    days = int(exp, 36)
    expires = ZP_EPOCH.timestamp() + days * 86400
    expires_dt = datetime.fromtimestamp(expires, tz=timezone.utc)
    days_left = (expires_dt.date() - datetime.now(timezone.utc).date()).days
    if days_left <= 0:
        return {"valid": False, "reason": f"Este código venció el {expires_dt.strftime('%d/%m/%Y')}."}
    return {"valid": True, "expires": expires_dt, "days_left": days_left}


# ── Estado ────────────────────────────────────────────────────────────────────

if "espacios" not in st.session_state:
    st.session_state.espacios = []
if "espacio_idx" not in st.session_state:
    st.session_state.espacio_idx = None
if "licencia_ok" not in st.session_state:
    st.session_state.licencia_ok = False


# ── Funciones ─────────────────────────────────────────────────────────────────

def ocr_imagen(datos: bytes) -> str:
    try:
        import pytesseract
        from PIL import Image
        img = Image.open(io.BytesIO(datos))
        texto = pytesseract.image_to_string(img, lang="eng+spa")
        return texto.strip()
    except Exception:
        return ""


def extraer_texto(datos: bytes, nombre: str) -> str:
    ext = Path(nombre).suffix.lower()
    if ext in (".txt", ".md"):
        return datos.decode("utf-8", errors="ignore")

    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tmp:
        tmp.write(datos)
        ruta = tmp.name
    try:
        if ext == ".pdf":
            from pypdf import PdfReader
            return "\n".join(p.extract_text() or "" for p in PdfReader(ruta).pages)
        if ext == ".docx":
            from docx import Document
            return "\n".join(p.text for p in Document(ruta).paragraphs)
        if ext == ".pptx":
            from pptx import Presentation
            textos = []
            for slide in Presentation(ruta).slides:
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        textos.append(shape.text_frame.text)
            return "\n".join(textos)
    finally:
        os.unlink(ruta)
    return datos.decode("utf-8", errors="ignore")


def clasificar_archivo(nombre: str) -> str:
    ext = Path(nombre).suffix.lower().lstrip(".")
    if ext in EXT_IMAGEN:
        return "imagen"
    return "documento"


def generar_feedback(
    api_key: str, proveedor: str, modelo_llm: str,
    nombre: str, guia: str, rubrica: str,
    contenido: str, tipo: str, veredicto: str,
    observaciones: str = "", pct_dudosas: float = 0.0,
) -> str:
    if proveedor == "OpenRouter":
        client = OpenAI(api_key=api_key, base_url="https://openrouter.ai/api/v1")
    else:
        client = OpenAI(api_key=api_key)

    bloque_guia = f"ACTIVITY GUIDE:\n---\n{guia}\n---" if guia.strip() else ""
    bloque_rubrica = f"EVALUATION RUBRIC:\n---\n{rubrica}\n---" if rubrica.strip() else ""
    veredicto_texto = VEREDICTO_APROBADO if veredicto == "Aprobado" else VEREDICTO_DEFICIENTE

    if tipo == "imagen":
        bloque_contenido = (
            f"STUDENT'S SUBMISSION (image):\n"
            f"---\n{contenido}\n---"
        )
    else:
        bloque_contenido = f"STUDENT'S SUBMISSION (text/document):\n---\n{contenido}\n---"

    bloque_obs = f"INSTRUCTOR'S ADDITIONAL NOTES:\n{observaciones}" if observaciones.strip() else ""

    prompt = PROMPT.format(
        nombre=nombre or "Student",
        bloque_guia=bloque_guia,
        bloque_rubrica=bloque_rubrica,
        veredicto_texto=veredicto_texto,
        bloque_contenido=bloque_contenido,
        bloque_observaciones=bloque_obs,
    )
    resp = client.chat.completions.create(
        model=modelo_llm,
        messages=[{"role": "user", "content": prompt}],
        temperature=0.4,
        max_tokens=1500,
    )
    return resp.choices[0].message.content


def _nombre_seguro(texto: str) -> str:
    import re
    limpio = re.sub(r"[^\w\s.-]", "", texto).strip()
    return re.sub(r"\s+", "_", limpio)[:80] or "entrega"


def zip_entregas(entregas: list[dict], nombre_espacio: str) -> bytes:
    buffer = io.BytesIO()
    with zipfile.ZipFile(buffer, "w", zipfile.ZIP_DEFLATED) as z:
        usados: set[str] = set()
        for e in entregas:
            if not e.get("retroalimentacion"):
                continue
            base = _nombre_seguro(e["estudiante"]) or "estudiante"
            nombre = base
            c = 1
            while nombre in usados:
                nombre = f"{base}_{c}"
                c += 1
            usados.add(nombre)
            contenido = (
                f"Estudiante: {e['estudiante']}\n"
                f"Archivo: {e['archivo_nombre']}\n"
                f"Veredicto: {e['veredicto']}\n\n"
                f"Retroalimentación:\n{e['retroalimentacion']}\n"
            )
            z.writestr(f"{nombre}.txt", contenido)
    return buffer.getvalue()


def excel_entregas(entregas: list[dict], nombre_espacio: str) -> bytes:
    filas = []
    for e in entregas:
        filas.append({
            "Estudiante": e["estudiante"],
            "Archivo": e["archivo_nombre"],
            "Tipo": e["tipo"],
            "Veredicto": e.get("veredicto", ""),
            "Retroalimentación": e.get("retroalimentacion", ""),
            "Observaciones": e.get("observaciones", ""),
        })
    buf = io.BytesIO()
    pd.DataFrame(filas).to_excel(buf, index=False, engine="openpyxl")
    return buf.getvalue()


# ── Gate de licencia ─────────────────────────────────────────────────────────

def mostrar_gate():
    st.title("🎓 RetroTutor")
    st.markdown(
        "Retroalimentación automática con IA para evidencias estudiantiles. "
        "Disponible para suscriptores de **Zajuna Panel Pro**."
    )
    st.divider()

    col1, col2 = st.columns([2, 1])
    with col1:
        st.subheader("Activa tu acceso")
        st.markdown(
            "Ingresa el código de licencia que recibiste con tu suscripción a Zajuna Panel. "
            "El mismo código que activa la extensión funciona aquí."
        )
        code = st.text_input(
            "Código de licencia", placeholder="ZP-XXX-XXXX-XXXXXX",
            help="Formato: ZP seguido de tres grupos separados por guiones",
        )
        if st.button("Activar", type="primary", disabled=not code):
            result = zp_validate(code)
            if result["valid"]:
                st.session_state.licencia_ok = True
                st.session_state.licencia_code = code
                st.session_state.licencia_days = result["days_left"]
                st.rerun()
            else:
                st.error(result["reason"])

    with col2:
        st.subheader("No tienes código?")
        st.markdown(
            "Adquiere **Zajuna Panel Pro** y obtén acceso a:\n"
            "- Dashboard de calificaciones\n"
            "- Portal del Aprendiz\n"
            "- **RetroTutor** (retroalimentación IA)\n"
            "- Exportación CSV\n\n"
            "Desde **$10.000 COP/mes**"
        )


if not st.session_state.licencia_ok:
    mostrar_gate()
    st.stop()


# ══════════════════════════════════════════════════════════════════════════════
# App principal (solo si la licencia es válida)
# ══════════════════════════════════════════════════════════════════════════════

# ── Barra lateral ────────────────────────────────────────────────────────────

with st.sidebar:
    st.header("⚙️ Configuración")

    days_left = st.session_state.get("licencia_days", 0)
    if days_left <= 7:
        st.warning(f"Tu licencia vence en {days_left} día(s).")
    else:
        st.success(f"Licencia activa ({days_left} días restantes)")

    st.subheader("Retroalimentación IA")
    proveedor = st.selectbox("Proveedor", ["OpenRouter", "OpenAI"])
    api_key = st.text_input(
        "API Key OpenRouter" if proveedor == "OpenRouter" else "API Key OpenAI",
        type="password",
        help="Obtén una gratis en openrouter.ai" if proveedor == "OpenRouter" else "platform.openai.com",
    )
    MODELOS = {
        "OpenRouter": [
            "nvidia/nemotron-3-super-120b-a12b:free",
            "deepseek/deepseek-r1-0528:free",
            "google/gemma-3-27b-it:free",
            "meta-llama/llama-4-maverick:free",
        ],
        "OpenAI": ["gpt-4o-mini", "gpt-4o", "gpt-3.5-turbo"],
    }
    modelo_llm = st.selectbox("Modelo LLM", MODELOS[proveedor], help="Los modelos :free no tienen costo en OpenRouter")


# ── Navegación ───────────────────────────────────────────────────────────────

idx = st.session_state.espacio_idx

if idx is None:
    # ── Vista HOME ───────────────────────────────────────────────────────
    st.title("🎓 RetroTutor")
    st.markdown(
        "**Retroalimentación automática para cualquier evidencia estudiantil:** "
        "documentos, texto, imágenes e infografías."
    )

    # Crear espacio
    with st.expander(
        "➕ Crear espacio de retroalimentación",
        expanded=len(st.session_state.espacios) == 0,
    ):
        nombre_ev = st.text_input(
            "Nombre de la evidencia",
            placeholder="Ej: Infografía Unidad 2, Presentación oral, Ensayo final…",
        )
        tab_file, tab_text = st.tabs(["📎 Subir archivos", "✏️ Escribir texto"])

        with tab_file:
            guia_file = st.file_uploader(
                "Guía de la actividad", type=EXT_DOCUMENTO, key="guia_f_crear",
            )
            rubrica_file = st.file_uploader(
                "Rúbrica de evaluación (opcional)", type=EXT_DOCUMENTO, key="rub_f_crear",
            )
        with tab_text:
            guia_txt = st.text_area(
                "Guía de la actividad", height=100, key="guia_txt_crear",
                placeholder="Describe la actividad y lo que se espera del estudiante…",
            )
            rubrica_txt = st.text_area(
                "Rúbrica (opcional)", height=100, key="rubrica_txt_crear",
                placeholder="Criterios de evaluación y niveles de desempeño…",
            )

        if st.button(
            "✅ Crear espacio", type="primary",
            use_container_width=True, disabled=not nombre_ev,
        ):
            guia = ""
            if guia_file:
                guia = extraer_texto(guia_file.read(), guia_file.name)
            elif guia_txt:
                guia = guia_txt
            rubrica = ""
            if rubrica_file:
                rubrica = extraer_texto(rubrica_file.read(), rubrica_file.name)
            elif rubrica_txt:
                rubrica = rubrica_txt

            st.session_state.espacios.append({
                "nombre": nombre_ev,
                "guia": guia,
                "rubrica": rubrica,
                "entregas": [],
            })
            st.session_state.espacio_idx = len(st.session_state.espacios) - 1
            st.rerun()

    # Lista de espacios existentes
    if st.session_state.espacios:
        st.divider()
        st.subheader("Tus espacios")
        for i, esp in enumerate(st.session_state.espacios):
            n_ent = len(esp["entregas"])
            n_retro = sum(1 for e in esp["entregas"] if e.get("retroalimentacion"))
            col_name, col_btn = st.columns([4, 1])
            with col_name:
                st.markdown(f"**{esp['nombre']}** — {n_ent} entregas, {n_retro} retroalimentadas")
            with col_btn:
                if st.button("Abrir", key=f"abrir_{i}"):
                    st.session_state.espacio_idx = i
                    st.rerun()

else:
    # ── Vista ESPACIO ────────────────────────────────────────────────────
    espacio = st.session_state.espacios[idx]

    if st.button("← Volver"):
        st.session_state.espacio_idx = None
        st.rerun()

    st.title(f"🎓 {espacio['nombre']}")

    has_guia = bool(espacio["guia"].strip())
    has_rubrica = bool(espacio["rubrica"].strip())
    st.caption(
        f"Guía: {'✅' if has_guia else '❌'} · "
        f"Rúbrica: {'✅' if has_rubrica else '—'} · "
        f"{len(espacio['entregas'])} entregas"
    )

    # ── Agregar entregas ─────────────────────────────────────────────────
    st.divider()
    st.subheader("Agregar entregas")

    tab_archivos, tab_texto_ev = st.tabs(["📎 Subir archivos", "✏️ Pegar texto"])

    with tab_archivos:
        archivos = st.file_uploader(
            "Subir entregas de estudiantes",
            type=TODAS_EXT,
            accept_multiple_files=True,
            key=f"upload_{idx}",
            help="PDF, DOCX, PPTX, TXT, imágenes con OCR (PNG, JPG). El nombre del archivo se usa como nombre del estudiante.",
        )
        if archivos and st.button(
            f"➕ Agregar {len(archivos)} entregas",
            type="primary", key=f"agregar_{idx}",
        ):
            barra = st.progress(0, text="Procesando entregas…")
            for i, arch in enumerate(archivos):
                barra.progress(i / len(archivos), text=f"Procesando {arch.name}…")
                tipo = clasificar_archivo(arch.name)
                raw = arch.read()
                contenido = ""

                if tipo == "documento":
                    contenido = extraer_texto(raw, arch.name)
                elif tipo == "imagen":
                    with st.spinner(f"Extrayendo texto de {arch.name}…"):
                        contenido = ocr_imagen(raw)

                espacio["entregas"].append({
                    "estudiante": Path(arch.name).stem.replace("_", " "),
                    "tipo": tipo,
                    "contenido": contenido,
                    "archivo_nombre": arch.name,
                    "veredicto": None,
                    "retroalimentacion": None,
                    "observaciones": "",
                    "imagen_bytes": raw if tipo == "imagen" else None,
                    "archivo_bytes": raw,
                    "archivo_mime": arch.type or "application/octet-stream",
                })
            barra.progress(1.0, text="✅ Entregas agregadas")
            st.rerun()

    with tab_texto_ev:
        nombre_txt = st.text_input(
            "Nombre del estudiante", key=f"nom_txt_{idx}",
            placeholder="Ej: María García",
        )
        texto_pegado = st.text_area(
            "Pegar la entrega del estudiante", height=200, key=f"txt_{idx}",
            placeholder="Pega aquí el texto que el estudiante envió…",
        )
        if st.button("➕ Agregar entrega", key=f"add_txt_{idx}") and texto_pegado:
            espacio["entregas"].append({
                "estudiante": nombre_txt or "Estudiante",
                "tipo": "texto",
                "contenido": texto_pegado,
                "archivo_nombre": "texto pegado",
                "veredicto": None,
                "retroalimentacion": None,
                "observaciones": "",
                "imagen_bytes": None,
                "archivo_bytes": None,
                "archivo_mime": None,
            })
            st.rerun()

    # ── Lista de entregas ────────────────────────────────────────────────
    if not espacio["entregas"]:
        st.info("Aún no hay entregas. Sube archivos o pega texto.")
    else:
        st.divider()

        sin_veredicto = sum(1 for e in espacio["entregas"] if not e.get("veredicto"))

        st.subheader(f"Calificar entregas ({len(espacio['entregas'])})")
        st.info(
            "**Paso 1:** Revisa cada entrega y marca **Aprobado** o **Deficiente**. "
            "**Paso 2:** Haz clic en **Generar retroalimentación**.",
            icon="👇",
        )

        for i, ent in enumerate(espacio["entregas"]):
            if ent.get("retroalimentacion"):
                icono = "✅"
            elif ent.get("veredicto"):
                icono = "⏳"
            else:
                icono = "⬜"

            with st.expander(
                f"{icono} {ent['estudiante']} — {ent['archivo_nombre']}",
                expanded=not ent.get("retroalimentacion"),
            ):
                col_prev, col_ctrl = st.columns([3, 2])

                with col_prev:
                    tipo_label = {
                        "documento": "📄 Documento",
                        "texto": "📝 Texto",
                        "imagen": "🖼️ Imagen",
                    }
                    st.caption(tipo_label.get(ent["tipo"], ent["tipo"]))

                    if ent.get("archivo_bytes"):
                        st.download_button(
                            f"📥 Descargar {ent['archivo_nombre']}",
                            data=ent["archivo_bytes"],
                            file_name=ent["archivo_nombre"],
                            mime=ent.get("archivo_mime", "application/octet-stream"),
                            key=f"dl_{idx}_{i}",
                        )

                    if ent["tipo"] == "imagen" and ent.get("imagen_bytes"):
                        st.image(ent["imagen_bytes"], width=400)
                        if ent.get("contenido"):
                            with st.container(border=True):
                                st.caption("Texto extraído (OCR)")
                                st.write(ent["contenido"][:3000])
                    else:
                        with st.container(border=True):
                            st.caption("Contenido extraído")
                            st.write(ent["contenido"][:3000] or "*Sin contenido*")

                with col_ctrl:
                    ent["estudiante"] = st.text_input(
                        "Nombre", value=ent["estudiante"], key=f"est_{idx}_{i}",
                    )

                    st.markdown("**Veredicto** *(requerido)*")
                    veredicto = st.radio(
                        "Veredicto", ["Aprobado", "Deficiente"],
                        index=(
                            0 if ent.get("veredicto") == "Aprobado"
                            else 1 if ent.get("veredicto") == "Deficiente"
                            else None
                        ),
                        key=f"ver_{idx}_{i}", horizontal=True,
                        label_visibility="collapsed",
                    )
                    ent["veredicto"] = veredicto

                    if ent["tipo"] == "imagen":
                        ent["observaciones"] = st.text_area(
                            "Describe lo que ves en la imagen",
                            value=ent.get("observaciones", ""),
                            key=f"obs_{idx}_{i}", height=80,
                            placeholder="Ej: Buen diseño visual, falta citar fuentes…",
                        )
                    else:
                        ent["observaciones"] = st.text_area(
                            "Observaciones del instructor (opcional)",
                            value=ent.get("observaciones", ""),
                            key=f"obs_{idx}_{i}", height=60,
                            placeholder="Notas adicionales…",
                        )

                    btn_label = "🔄 Regenerar" if ent.get("retroalimentacion") else "🚀 Generar retroalimentación"
                    if st.button(
                        btn_label, key=f"gen1_{idx}_{i}",
                        use_container_width=True,
                        disabled=not ent.get("veredicto"),
                        type="primary" if ent.get("veredicto") and not ent.get("retroalimentacion") else "secondary",
                    ):
                        if not api_key:
                            st.error("Ingresa tu API Key en la barra lateral.")
                        else:
                            contenido_llm = ent["contenido"]
                            if ent["tipo"] == "imagen":
                                partes = []
                                if ent.get("contenido"):
                                    partes.append(f"[OCR text from image]: {ent['contenido']}")
                                if ent.get("observaciones"):
                                    partes.append(f"[Instructor notes]: {ent['observaciones']}")
                                contenido_llm = "\n".join(partes) or "No description provided."
                            with st.spinner("Generando…"):
                                ent["retroalimentacion"] = generar_feedback(
                                    api_key, proveedor, modelo_llm,
                                    ent["estudiante"], espacio["guia"],
                                    espacio["rubrica"], contenido_llm,
                                    ent["tipo"], ent["veredicto"],
                                    ent.get("observaciones", ""),
                                )
                            st.rerun()

                    if not ent.get("veredicto"):
                        st.caption("↑ Selecciona Aprobado o Deficiente para habilitar")

                if ent.get("retroalimentacion"):
                    st.divider()
                    st.markdown("**Retroalimentación generada:**")
                    st.markdown(ent["retroalimentacion"])

                if st.button("🗑️ Eliminar entrega", key=f"del_{idx}_{i}"):
                    espacio["entregas"].pop(i)
                    st.rerun()

        # ── Acciones en lote ─────────────────────────────────────────────
        st.divider()
        pendientes = [
            e for e in espacio["entregas"]
            if e.get("veredicto") and not e.get("retroalimentacion")
        ]
        con_retro = [e for e in espacio["entregas"] if e.get("retroalimentacion")]

        if sin_veredicto > 0:
            st.warning(
                f"**{sin_veredicto} entrega(s) sin veredicto.** "
                "Marca Aprobado o Deficiente en cada una para poder generar retroalimentación.",
                icon="⚠️",
            )

        col_gen, col_zip, col_xls = st.columns(3)

        with col_gen:
            btn_disabled = len(pendientes) == 0 or not api_key
            if st.button(
                f"🚀 Generar {len(pendientes)} pendientes",
                type="primary", use_container_width=True,
                disabled=btn_disabled,
            ):
                barra = st.progress(0, text="Generando retroalimentaciones…")
                for j, e in enumerate(pendientes):
                    barra.progress(
                        j / len(pendientes),
                        text=f"Retroalimentando {e['estudiante']}…",
                    )
                    contenido_llm = e["contenido"]
                    if e["tipo"] == "imagen":
                        partes = []
                        if e.get("contenido"):
                            partes.append(f"[OCR text from image]: {e['contenido']}")
                        if e.get("observaciones"):
                            partes.append(f"[Instructor notes]: {e['observaciones']}")
                        contenido_llm = "\n".join(partes) or "No description provided."
                    e["retroalimentacion"] = generar_feedback(
                        api_key, proveedor, modelo_llm,
                        e["estudiante"], espacio["guia"], espacio["rubrica"],
                        contenido_llm, e["tipo"], e["veredicto"],
                        e.get("observaciones", ""),
                    )
                barra.progress(1.0, text="✅ Listo")
                st.rerun()

        with col_zip:
            if con_retro:
                st.download_button(
                    f"📦 Descargar .zip ({len(con_retro)})",
                    zip_entregas(con_retro, espacio["nombre"]),
                    file_name=f"retro_{_nombre_seguro(espacio['nombre'])}.zip",
                    mime="application/zip",
                    use_container_width=True,
                )

        with col_xls:
            if con_retro:
                st.download_button(
                    f"📊 Descargar .xlsx ({len(con_retro)})",
                    excel_entregas(espacio["entregas"], espacio["nombre"]),
                    file_name=f"retro_{_nombre_seguro(espacio['nombre'])}.xlsx",
                    mime="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
                    use_container_width=True,
                )

        if not api_key and pendientes:
            st.warning("Ingresa tu API Key en la barra lateral para generar retroalimentaciones.")
